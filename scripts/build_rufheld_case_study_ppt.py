"""
Build Rufheld case study PowerPoint with embedded stock images (Unsplash).
Run: python scripts/build_rufheld_case_study_ppt.py
"""
from __future__ import annotations

import sys
from pathlib import Path

import requests
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

# Unsplash — free to use; we download and embed (no hotlinking in deck).
IMAGES = [
    (
        "hero",
        "https://images.unsplash.com/photo-1556761175-5973dc0f32e7?w=1600&q=80",
    ),
    (
        "challenge",
        "https://images.unsplash.com/photo-1454165804606-c3d57bc86b40?w=1600&q=80",
    ),
    (
        "automation",
        "https://images.unsplash.com/photo-1551434678-e076c223a692?w=1600&q=80",
    ),
    (
        "results",
        "https://images.unsplash.com/photo-1522071820081-009f0129c71c?w=1600&q=80",
    ),
]


def download_images(dest_dir: Path) -> dict[str, Path]:
    dest_dir.mkdir(parents=True, exist_ok=True)
    out: dict[str, Path] = {}
    for name, url in IMAGES:
        path = dest_dir / f"{name}.jpg"
        r = requests.get(url, timeout=60)
        r.raise_for_status()
        path.write_bytes(r.content)
        out[name] = path
    return out


def set_title(tf, text: str, size_pt: int = 36, color: RGBColor | None = None):
    tf.clear()
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size_pt)
    p.font.bold = True
    if color:
        p.font.color.rgb = color


def add_body(slide, left, top, width, height, lines: list[str], size_pt: int = 18):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        p.font.size = Pt(size_pt)
        p.space_after = Pt(6)
        p.level = 0


def main() -> int:
    base = Path(__file__).resolve().parent.parent
    img_dir = base / "storage" / "app" / "rufheld_ppt_assets"
    try:
        imgs = download_images(img_dir)
    except Exception as e:
        print("Image download failed:", e, file=sys.stderr)
        return 1

    prs = Presentation()
    prs.slide_width = Inches(13.333)  # 16:9
    prs.slide_height = Inches(7.5)

    blank = prs.slide_layouts[6]

    # --- Slide 1: Cover ---
    s1 = prs.slides.add_slide(blank)
    s1.shapes.add_picture(str(imgs["hero"]), 0, 0, width=prs.slide_width, height=prs.slide_height)
    # Dark overlay bar for text readability
    bar = s1.shapes.add_shape(1, Inches(0), Inches(2.2), prs.slide_width, Inches(3.2))  # rectangle
    bar.fill.solid()
    bar.fill.fore_color.rgb = RGBColor(15, 23, 42)
    bar.fill.transparency = 0.35
    bar.line.fill.background()
    tb = s1.shapes.add_textbox(Inches(0.6), Inches(2.5), Inches(12), Inches(2.8))
    tf = tb.text_frame
    tf.clear()
    for i, (t, sz) in enumerate(
        [
            ("SUCCESS STORY", 14),
            ("Rufheld", 44),
            ("Streamlining Reputation Management", 22),
            ("XPERT PRIME CASE STUDY", 12),
        ]
    ):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = t
        p.font.size = Pt(sz)
        p.font.bold = i in (1, 2)
        p.font.color.rgb = RGBColor(248, 250, 252)
        p.alignment = PP_ALIGN.LEFT
        if i == 0:
            p.space_after = Pt(4)

    # --- Slide 2: Who we help + image right ---
    s2 = prs.slides.add_slide(blank)
    s2.background.fill.solid()
    s2.background.fill.fore_color.rgb = RGBColor(248, 250, 252)
    s2.shapes.add_picture(str(imgs["automation"]), Inches(7.2), Inches(0.4), width=Inches(5.9))
    tbox = s2.shapes.add_textbox(Inches(0.55), Inches(0.45), Inches(6.2), Inches(1))
    set_title(tbox.text_frame, "Who We Help", 32, RGBColor(15, 23, 42))
    add_body(
        s2,
        Inches(0.55),
        Inches(1.35),
        Inches(6.2),
        Inches(5.6),
        [
            "INDUSTRY: Digital services · Online reputation management",
            "",
            "Rufheld is Germany-focused: removal of unfair or policy-violating Google reviews.",
            "Protects online reputation through expert analysis and strategic removal requests.",
            "",
            "• Pay only on success — no upfront cost",
            "• Fast turnaround — often within days",
            "• Manual analysis — higher success rates",
        ],
        16,
    )

    # --- Slide 3: Challenge ---
    s3 = prs.slides.add_slide(blank)
    s3.shapes.add_picture(str(imgs["challenge"]), Inches(0), Inches(0), width=Inches(5.2), height=prs.slide_height)
    tbox = s3.shapes.add_textbox(Inches(5.6), Inches(0.5), Inches(7.2), Inches(0.9))
    set_title(tbox.text_frame, "The Challenge", 32, RGBColor(15, 23, 42))
    add_body(
        s3,
        Inches(5.6),
        Inches(1.45),
        Inches(7),
        Inches(5.5),
        [
            "No centralized platform — scattered client data and cases",
            "Manual tracking — inefficiencies and human error",
            "Manual payment collection — delays in revenue recognition",
            "No case-to-billing connection — revenue leakage and misalignment",
        ],
        18,
    )

    # --- Slide 4: Strategic approach ---
    s4 = prs.slides.add_slide(blank)
    s4.background.fill.solid()
    s4.background.fill.fore_color.rgb = RGBColor(30, 41, 59)
    tbox = s4.shapes.add_textbox(Inches(0.55), Inches(0.45), Inches(12), Inches(0.85))
    set_title(tbox.text_frame, "Our Strategic Approach", 32, RGBColor(248, 250, 252))
    add_body(
        s4,
        Inches(0.55),
        Inches(1.35),
        Inches(12.2),
        Inches(5.8),
        [
            "Automated workflows — lead to payment, end-to-end",
            "Customized CRM architecture — tailored to success-based model",
            "Structured pipeline tracking — full lifecycle per case",
            "Automation-driven workflows — fewer manual steps and errors",
            "Trigger-based Stripe integration — bill on successful outcome",
            "Validation layers — data accuracy and process control",
            "Scalability and performance — built for growth",
        ],
        17,
    )
    # Body text on dark slide: light gray
    for shape in s4.shapes:
        if not shape.has_text_frame:
            continue
        if shape is tbox:
            continue
        for para in shape.text_frame.paragraphs:
            para.font.color.rgb = RGBColor(226, 232, 240)

    # --- Slide 5: Solution ---
    s5 = prs.slides.add_slide(blank)
    s5.background.fill.solid()
    s5.background.fill.fore_color.rgb = RGBColor(255, 255, 255)
    s5.shapes.add_picture(str(imgs["automation"]), Inches(7), Inches(1.1), width=Inches(6))
    tbox = s5.shapes.add_textbox(Inches(0.55), Inches(0.45), Inches(6.3), Inches(0.85))
    set_title(tbox.text_frame, "The Solution in Action", 30, RGBColor(15, 23, 42))
    add_body(
        s5,
        Inches(0.55),
        Inches(1.25),
        Inches(6.2),
        Inches(5.8),
        [
            "Zoho CRM: Leads, Accounts, Contacts, Review Cases",
            "Lifecycle: New → Request → Analysis submitted → In progress → Successful / Failed",
            "Stripe: success-based billing; payment after success; minimal manual billing steps",
            "Deluge + scripts + UI validations — backend logic and controls",
            "Real-time dashboard — operations and performance metrics",
        ],
        15,
    )

    # --- Slide 6: Ecosystem ---
    s6 = prs.slides.add_slide(blank)
    s6.background.fill.solid()
    s6.background.fill.fore_color.rgb = RGBColor(241, 245, 249)
    tbox = s6.shapes.add_textbox(Inches(0.55), Inches(0.45), Inches(12), Inches(0.85))
    set_title(tbox.text_frame, "Powered by Zoho Ecosystem", 30, RGBColor(15, 23, 42))
    cols = [
        ("Zoho CRM", "Core data and automation"),
        ("Zoho Books", "Payment management"),
        ("Zoho Flow", "Stripe integration"),
        ("Zoho Analytics", "Real-time insights"),
        ("Stripe", "Payment automation"),
        ("Deluge + APIs", "Logic, CRM–Stripe, UI validations"),
    ]
    x0, y0, cw, ch, gap = 0.55, 1.35, 3.85, 1.55, 0.35
    for i, (h, sub) in enumerate(cols):
        row, col = divmod(i, 3)
        left = Inches(x0 + col * (cw + gap))
        top = Inches(y0 + row * (ch + gap))
        sh = s6.shapes.add_shape(1, left, top, Inches(cw), Inches(ch))
        sh.fill.solid()
        sh.fill.fore_color.rgb = RGBColor(255, 255, 255)
        sh.line.color.rgb = RGBColor(203, 213, 225)
        tf = sh.text_frame
        tf.margin_left = Inches(0.12)
        tf.margin_top = Inches(0.12)
        p = tf.paragraphs[0]
        p.text = h
        p.font.bold = True
        p.font.size = Pt(16)
        p.font.color.rgb = RGBColor(15, 23, 42)
        p2 = tf.add_paragraph()
        p2.text = sub
        p2.font.size = Pt(12)
        p2.font.color.rgb = RGBColor(71, 85, 105)

    # --- Slide 7: Results + image ---
    s7 = prs.slides.add_slide(blank)
    s7.shapes.add_picture(str(imgs["results"]), Inches(6.8), Inches(0.35), width=Inches(6.35))
    tbox = s7.shapes.add_textbox(Inches(0.55), Inches(0.45), Inches(6), Inches(0.85))
    set_title(tbox.text_frame, "Transformative Results", 30, RGBColor(15, 23, 42))
    add_body(
        s7,
        Inches(0.55),
        Inches(1.25),
        Inches(6.1),
        Inches(5.9),
        [
            "Manual process reduction — automated lead-to-payment path",
            "Faster payment collection — automated billing and tracking",
            "Complete visibility — real-time operations and revenue insight",
            "Scalable growth — infrastructure for higher request volume",
            "Seamless flow — inquiry through successful case and payment",
            "Business intelligence — success rates, revenue, performance metrics",
        ],
        15,
    )

    # --- Slide 8: Closing ---
    s8 = prs.slides.add_slide(blank)
    s8.shapes.add_picture(str(imgs["hero"]), 0, 0, width=prs.slide_width, height=prs.slide_height)
    bar = s8.shapes.add_shape(1, Inches(0), Inches(2.4), prs.slide_width, Inches(2.9))
    bar.fill.solid()
    bar.fill.fore_color.rgb = RGBColor(15, 23, 42)
    bar.fill.transparency = 0.25
    bar.line.fill.background()
    tb = s8.shapes.add_textbox(Inches(0.7), Inches(2.55), Inches(11.8), Inches(2.5))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "A Complete Transformation"
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)
    p2 = tf.add_paragraph()
    p2.text = (
        "Rufheld now runs a centralized, automated, scalable system aligned with "
        "success-based delivery — from CRM and case lifecycle to Stripe billing and real-time dashboards."
    )
    p2.font.size = Pt(15)
    p2.font.color.rgb = RGBColor(226, 232, 240)

    out = Path.home() / "Downloads" / "Rufheld_Streamlining_Reputation_Management.pptx"
    prs.save(str(out))
    print("Saved:", out)
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
